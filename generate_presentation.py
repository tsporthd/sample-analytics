#!/usr/bin/env python3
"""
Generate PowerPoint presentation explaining the risk analysis logic.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    """Create a PowerPoint presentation explaining the analytics logic."""

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Overview
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Title
    title_box = slide1.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Infrastructure Risk Analysis - Logic Overview"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 51, 102)
    title_para.alignment = PP_ALIGN.CENTER

    # Content box
    content_box = slide1.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    # Data Processing Flow
    p = tf.paragraphs[0]
    p.text = "Data Processing Flow:"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.space_after = Pt(12)

    flow_items = [
        "1. Load Infrastructure Data: Read from sampleData.csv",
        "2. Filter by AppCode: Apply Apps.csv filter to include only in-scope applications",
        "3. Deduplicate: Keep first occurrence of each unique AppCode",
        "4. Count Infrastructure: Calculate total infrastructure items per AppCode",
        "5. Calculate Risk Metrics: Apply scoring formulas to each record",
        "6. Generate Reports: Create risk chart ranked by composite score and risk percentage"
    ]

    for item in flow_items:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.space_after = Pt(8)
        p.level = 0

    # Key Principle
    p = tf.add_paragraph()
    p.text = "\nKey Principle:"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(204, 0, 0)
    p.space_before = Pt(16)

    p = tf.add_paragraph()
    p.text = "Higher composite scores + more infrastructure items = Higher overall risk"
    p.font.size = Pt(16)
    p.font.italic = True
    p.font.color.rgb = RGBColor(102, 102, 102)

    # Slide 2: Risk Calculation Formulas
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Title
    title_box2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame2 = title_box2.text_frame
    title_frame2.text = "Risk Calculation & Ranking Methodology"
    title_para2 = title_frame2.paragraphs[0]
    title_para2.font.size = Pt(32)
    title_para2.font.bold = True
    title_para2.font.color.rgb = RGBColor(0, 51, 102)
    title_para2.alignment = PP_ALIGN.CENTER

    # Left column - Formulas
    left_box = slide2.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(4.5), Inches(5.5))
    tf_left = left_box.text_frame
    tf_left.word_wrap = True

    p = tf_left.paragraphs[0]
    p.text = "Calculation Formulas:"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.space_after = Pt(12)

    # CompositeScoreNumber mapping
    p = tf_left.add_paragraph()
    p.text = "1. CompositeScoreNumber:"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_before = Pt(8)

    score_mappings = [
        "   • High = 3.0",
        "   • Moderate High = 2.5",
        "   • Moderate = 2.0",
        "   • Low = 1.0"
    ]

    for mapping in score_mappings:
        p = tf_left.add_paragraph()
        p.text = mapping
        p.font.size = Pt(14)
        p.space_after = Pt(4)

    # Other formulas
    formulas = [
        ("2. TotalInfrastructure:", "Count of infrastructure items per AppCode"),
        ("3. CompositeRiskScore:", "CompositeScoreNumber × TotalInfrastructure"),
        ("4. CompositeRiskScorePercent:", "(CompositeRiskScore / TotalRiskScore) × 100")
    ]

    for title, formula in formulas:
        p = tf_left.add_paragraph()
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.space_before = Pt(12)

        p = tf_left.add_paragraph()
        p.text = f"   {formula}"
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(51, 51, 51)

    # Right column - Ranking
    right_box = slide2.shapes.add_textbox(Inches(5.2), Inches(1.3), Inches(4.3), Inches(5.5))
    tf_right = right_box.text_frame
    tf_right.word_wrap = True

    p = tf_right.paragraphs[0]
    p.text = "Ranking Methodology:"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.space_after = Pt(12)

    p = tf_right.add_paragraph()
    p.text = "Primary Sort:"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_before = Pt(8)

    p = tf_right.add_paragraph()
    p.text = "CompositeScoreNumber (High to Low)"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(51, 51, 51)
    p.space_after = Pt(8)

    p = tf_right.add_paragraph()
    p.text = "Secondary Sort:"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_before = Pt(12)

    p = tf_right.add_paragraph()
    p.text = "CompositeRiskScorePercent (High to Low)"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(51, 51, 51)
    p.space_after = Pt(12)

    # Example
    p = tf_right.add_paragraph()
    p.text = "Example Calculation:"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(204, 0, 0)
    p.space_before = Pt(16)

    example_lines = [
        "AppCode: CATIOT",
        "CompositeScore: Moderate (2.0)",
        "TotalInfrastructure: 1,237 items",
        "→ RiskScore: 2.0 × 1,237 = 2,474",
        "→ Percent: 2,474 / 12,932 = 19.1%"
    ]

    for line in example_lines:
        p = tf_right.add_paragraph()
        p.text = line
        p.font.size = Pt(13)
        p.space_after = Pt(4)
        if "→" in line:
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 102, 0)

    # Save presentation
    prs.save('Infrastructure_Risk_Analysis.pptx')
    print("PowerPoint presentation created: Infrastructure_Risk_Analysis.pptx")

if __name__ == "__main__":
    create_presentation()
